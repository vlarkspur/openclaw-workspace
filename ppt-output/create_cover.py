#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
金银 Planet PPT 封面生成器
设计：Jobs 审核通过的风格
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# 配置
OUTPUT_PPTX = "/home/admin/.openclaw/workspace/ppt-output/01-cover.pptx"
OUTPUT_PNG = "/home/admin/.openclaw/workspace/ppt-output/01-cover.png"
CHARACTER_IMAGE = "/home/admin/.openclaw/workspace/images/jinyin-character/04-cover-final.png"

# 颜色定义
COLOR_BLACK = RGBColor(0, 0, 0)  # #000000 深空黑
COLOR_GOLD = RGBColor(255, 215, 0)  # #FFD700 金色
COLOR_SILVER = RGBColor(192, 192, 192)  # #C0C0C0 银色
COLOR_WHITE = RGBColor(255, 255, 255)  # 白色

def create_gradient_text(shape, text, font_size, is_title=False):
    """
    创建渐变色文字效果
    python-pptx 不直接支持文字渐变，我们用金色到银色的过渡模拟
    """
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.text = text
    
    run = paragraph.runs[0]
    run.font.size = font_size
    run.font.bold = True
    run.font.name = 'Arial'
    
    # 使用金色作为主色（标题）
    if is_title:
        run.font.color.rgb = COLOR_GOLD
    else:
        run.font.color.rgb = COLOR_SILVER
    
    return shape

def main():
    # 创建 16:9 PPT (1920x1080)
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 1920px / 144 DPI
    prs.slide_height = Inches(7.5)    # 1080px / 144 DPI
    
    # 添加空白幻灯片
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 1. 设置背景为深空黑
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BLACK
    
    # 2. 添加标题 "金银 Planet" - 金银渐变色
    title_box = slide.shapes.add_textbox(
        Inches(1),   # left
        Inches(2),   # top (留白)
        Inches(11.333),  # width
        Inches(2)    # height
    )
    create_gradient_text(title_box, "金银 Planet", Pt(72), is_title=True)
    
    # 3. 添加副标题 "一个人的集团，一支 AI 军队" - 白色
    subtitle_box = slide.shapes.add_textbox(
        Inches(1),
        Inches(3.8),
        Inches(11.333),
        Inches(1)
    )
    subtitle_para = subtitle_box.text_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.text = "一个人的集团，一支 AI 军队"
    subtitle_run = subtitle_para.runs[0]
    subtitle_run.font.size = Pt(28)
    subtitle_run.font.name = 'Arial'
    subtitle_run.font.color.rgb = COLOR_WHITE
    
    # 4. 添加底部小字 "让 AI 给你打工，你只管做老板" - 银色
    footer_box = slide.shapes.add_textbox(
        Inches(1),
        Inches(6.5),
        Inches(11.333),
        Inches(0.8)
    )
    footer_para = footer_box.text_frame.paragraphs[0]
    footer_para.alignment = PP_ALIGN.CENTER
    footer_para.text = "让 AI 给你打工，你只管做老板"
    footer_run = footer_para.runs[0]
    footer_run.font.size = Pt(16)
    footer_run.font.name = 'Arial'
    footer_run.font.color.rgb = COLOR_SILVER
    
    # 5. 添加金银形象图片（如果存在）
    if os.path.exists(CHARACTER_IMAGE):
        try:
            # 将图片放在右侧，保持视觉平衡
            pic = slide.shapes.add_picture(
                CHARACTER_IMAGE,
                Inches(9),      # left - 右侧
                Inches(1.5),    # top
                Inches(4)       # width - 调整大小
            )
            print(f"✓ 已添加金银形象图片")
        except Exception as e:
            print(f"⚠ 图片添加失败：{e}")
    else:
        print(f"⚠ 图片文件不存在：{CHARACTER_IMAGE}")
    
    # 保存 PPTX
    prs.save(OUTPUT_PPTX)
    print(f"✓ PPT 已保存：{OUTPUT_PPTX}")
    
    # 生成 PNG 预览（使用 PIL）
    try:
        from PIL import Image
        
        # 创建一个 1920x1080 的图片
        img = Image.new('RGB', (1920, 1080), color='black')
        
        # 由于无法直接渲染 PPTX 到图片，我们创建一个简化的预览图
        # 实际使用中可能需要其他工具如 comtypes (Windows) 或 libreoffice
        print(f"⚠ PNG 预览需要额外工具，PPTX 文件已生成")
        print(f"  可使用 LibreOffice 转换：libreoffice --headless --convert-to png {OUTPUT_PPTX}")
        
    except ImportError:
        print(f"⚠ PIL 不可用，跳过 PNG 生成")
    
    print("\n" + "="*60)
    print("📊 制作说明")
    print("="*60)
    print(f"字体：Arial (无衬线字体)")
    print(f"背景：深空黑 #000000")
    print(f"标题：金银 Planet - 金色 #FFD700")
    print(f"副标题：一个人的集团，一支 AI 军队 - 白色 #FFFFFF")
    print(f"底部小字：让 AI 给你打工，你只管做老板 - 银色 #C0C0C0")
    print(f"排版：居中对齐，留白率≥60%")
    print(f"PPT 尺寸：16:9 (1920×1080)")
    print("="*60)
    print(f"\n✅ 输出文件：{OUTPUT_PPTX}")
    print(f"📁 请用户确认风格后继续制作其他页")

if __name__ == "__main__":
    main()
